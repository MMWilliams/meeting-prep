#!/usr/bin/env python3
"""
Meeting Prep CLI Tool
A comprehensive tool for engineers to prepare for meetings by analyzing documents
and generating detailed technical briefings.
"""

import argparse
import os
import re
import json
import datetime
from pathlib import Path
from typing import List, Dict, Optional
import sys

# Core dependencies
import PyPDF2
from PIL import Image
import pytesseract
from pptx import Presentation
import docx
import chardet
import spacy
from presidio_analyzer import AnalyzerEngine
from presidio_anonymizer import AnonymizerEngine
import openai
from reportlab.lib.pagesizes import letter
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, PageBreak
from reportlab.lib import colors
from reportlab.lib.enums import TA_CENTER, TA_JUSTIFY
from rich.console import Console
from rich.prompt import Prompt
from rich.progress import Progress, SpinnerColumn, TextColumn
from tenacity import retry, stop_after_attempt, wait_exponential
from dotenv import load_dotenv


console = Console()


class ContentExtractor:
    """Extracts text content from various file formats"""
    
    def extract(self, file_path: Path) -> str:
        suffix = file_path.suffix.lower()
        
        extractors = {
            '.pdf': self.extract_pdf,
            '.txt': self.extract_text,
            '.png': self.extract_image,
            '.jpg': self.extract_image,
            '.jpeg': self.extract_image,
            '.pptx': self.extract_pptx,
            '.docx': self.extract_docx
        }
        
        extractor = extractors.get(suffix)
        if not extractor:
            raise ValueError(f"Unsupported file type: {suffix}")
            
        return extractor(file_path)
    
    def extract_pdf(self, file_path: Path) -> str:
        text = ""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page in pdf_reader.pages:
                    text += page.extract_text() + "\n"
        except Exception as e:
            console.print(f"[yellow]Warning: Error extracting PDF {file_path.name}: {e}[/yellow]")
        return text
    
    def extract_text(self, file_path: Path) -> str:
        try:
            with open(file_path, 'rb') as file:
                raw_data = file.read()
                result = chardet.detect(raw_data)
                encoding = result['encoding'] or 'utf-8'
            
            with open(file_path, 'r', encoding=encoding) as file:
                return file.read()
        except Exception as e:
            console.print(f"[yellow]Warning: Error reading text file {file_path.name}: {e}[/yellow]")
            return ""
    
    def extract_image(self, file_path: Path) -> str:
        try:
            image = Image.open(file_path)
            return pytesseract.image_to_string(image)
        except Exception as e:
            console.print(f"[yellow]Warning: Error extracting text from image {file_path.name}: {e}[/yellow]")
            return ""
    
    def extract_pptx(self, file_path: Path) -> str:
        text = ""
        try:
            prs = Presentation(str(file_path))
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            console.print(f"[yellow]Warning: Error extracting PowerPoint {file_path.name}: {e}[/yellow]")
        return text
    
    def extract_docx(self, file_path: Path) -> str:
        try:
            doc = docx.Document(str(file_path))
            return "\n".join([paragraph.text for paragraph in doc.paragraphs])
        except Exception as e:
            console.print(f"[yellow]Warning: Error extracting Word document {file_path.name}: {e}[/yellow]")
            return ""


class SensitiveDataRemover:
    """Removes sensitive information from text using multiple methods"""
    
    def __init__(self):
        # Load spaCy model for NER
        try:
            self.nlp = spacy.load("en_core_web_sm")
        except OSError:
            console.print("[yellow]SpaCy model not found. Please run: python -m spacy download en_core_web_sm[/yellow]")
            sys.exit(1)
        
        # Initialize Presidio for PII detection
        self.analyzer = AnalyzerEngine()
        self.anonymizer = AnonymizerEngine()
        
        # Custom patterns for additional sensitive data
        self.patterns = {
            'email': r'\b[A-Za-z0-9._%+-]+@[A-Za-z0-9.-]+\.[A-Z|a-z]{2,}\b',
            'phone': r'\b\d{3}[-.]?\d{3}[-.]?\d{4}\b',
            'ssn': r'\b\d{3}-\d{2}-\d{4}\b',
            'credit_card': r'\b\d{4}[-\s]?\d{4}[-\s]?\d{4}[-\s]?\d{4}\b',
            'ip_address': r'\b(?:\d{1,3}\.){3}\d{1,3}\b',
            'api_key': r'\b[A-Za-z0-9_]{20,}\b'
        }
    
    def clean_text(self, text: str) -> str:
        if not text:
            return ""
            
        try:
            # Use Presidio for comprehensive PII detection
            results = self.analyzer.analyze(text=text, language='en')
            anonymized_result = self.anonymizer.anonymize(text=text, analyzer_results=results)
            cleaned_text = anonymized_result.text
            
            # Additional pattern-based cleaning
            for pattern_name, pattern in self.patterns.items():
                cleaned_text = re.sub(pattern, f'[{pattern_name.upper()}_REDACTED]', cleaned_text)
            
            # Use spaCy for additional NER-based cleaning
            doc = self.nlp(cleaned_text[:1000000])  # Limit to 1M chars for performance
            for ent in doc.ents:
                if ent.label_ in ['PERSON', 'ORG', 'GPE']:
                    cleaned_text = cleaned_text.replace(ent.text, f'[{ent.label_}_REDACTED]')
            
            return cleaned_text
        except Exception as e:
            console.print(f"[yellow]Warning: Error cleaning sensitive data: {e}[/yellow]")
            return text


class DocumentProcessor:
    """Processes documents and extracts cleaned content"""
    
    def __init__(self, document_path: Path, verbose: bool = False):
        self.document_path = document_path
        self.verbose = verbose
        self.sensitive_data_remover = SensitiveDataRemover()
        self.content_extractor = ContentExtractor()
        self.supported_extensions = ['.pdf', '.txt', '.png', '.jpg', '.jpeg', '.pptx', '.docx']
    
    def process_all(self) -> List[Dict[str, str]]:
        documents = []
        
        if not self.document_path.exists():
            raise FileNotFoundError(f"Path does not exist: {self.document_path}")
            
        files = list(self.document_path.glob('**/*'))
        supported_files = [f for f in files if f.suffix.lower() in self.supported_extensions]
        
        if not supported_files:
            console.print(f"[yellow]No supported files found in {self.document_path}[/yellow]")
            return documents
            
        for file_path in supported_files:
            if self.verbose:
                console.print(f"Processing: {file_path.name}")
                
            content = self.content_extractor.extract(file_path)
            if content:
                cleaned_content = self.sensitive_data_remover.clean_text(content)
                documents.append({
                    'filename': file_path.name,
                    'content': cleaned_content,
                    'type': file_path.suffix
                })
        
        return documents


class OpenAIManager:
    """Manages OpenAI API interactions"""
    
    def __init__(self):
        self.api_key = os.getenv('OPENAI_API_KEY')
        if not self.api_key:
            raise ValueError("OPENAI_API_KEY environment variable not set")
        openai.api_key = self.api_key
        self.client = openai.OpenAI()
    
    @retry(stop=stop_after_attempt(3), wait=wait_exponential(multiplier=1, min=4, max=10))
    def generate_report(self, content: str, topic: str = None) -> Dict:
        if topic:
            prompt = f"Create a comprehensive technical briefing about: {topic}"
        else:
            prompt = f"Based on the following content, create a comprehensive technical briefing:\n\n{content[:10000]}"  # Limit content length
        
        prompt += """
        
        Please provide a detailed report with the following sections:
        1. Executive Summary
        2. Topic Overview
        3. Technology Stack Analysis
        4. Architecture Overview (if applicable)
        5. Advantages and Benefits
        6. Limitations and Challenges
        7. Alternative Solutions
        8. Competitive Analysis
        9. Key Recommendations
        10. Action Items and Next Steps
        
        Format the response as valid JSON with these sections as keys. Each section should be a string or array of strings.
        """
        
        try:
            response = self.client.chat.completions.create(
                model="gpt-4",  # Changed from gpt-4-turbo-preview to gpt-4
                messages=[
                    {"role": "system", "content": "You are a technical analyst preparing briefing documents for engineering meetings. Always respond with valid JSON."},
                    {"role": "user", "content": prompt}
                ],
                temperature=0.7
            )
            
            content = response.choices[0].message.content
            # Clean up content to ensure valid JSON
            content = content.strip()
            if content.startswith("```json"):
                content = content[7:]
            if content.endswith("```"):
                content = content[:-3]
            
            return json.loads(content)
        except json.JSONDecodeError:
            # Fallback structure if JSON parsing fails
            return self._create_fallback_report(topic or "Technical Briefing")
        except Exception as e:
            console.print(f"[red]Error calling OpenAI API: {e}[/red]")
            return self._create_fallback_report(topic or "Technical Briefing")
    
    def _create_fallback_report(self, topic: str) -> Dict:
        """Create a fallback report structure"""
        return {
            "Executive Summary": f"Unable to generate detailed analysis for {topic}. Please review source documents manually.",
            "Topic Overview": "Analysis pending manual review.",
            "Technology Stack Analysis": "Not available.",
            "Architecture Overview": "Not available.",
            "Advantages and Benefits": ["Requires manual analysis"],
            "Limitations and Challenges": ["Requires manual analysis"],
            "Alternative Solutions": ["Requires further research"],
            "Competitive Analysis": "Not available.",
            "Key Recommendations": ["Review source documents", "Conduct manual analysis"],
            "Action Items and Next Steps": ["Schedule follow-up discussion", "Gather additional information"]
        }


class ReportGenerator:
    """Generates reports from content or topics"""
    
    def __init__(self):
        self.openai_manager = OpenAIManager()
    
    def generate_from_content(self, documents: List[Dict[str, str]]) -> Dict:
        if not documents:
            return self._create_empty_report()
            
        # Combine all document content
        combined_content = "\n\n".join([
            f"Document: {doc['filename']}\nContent: {doc['content'][:2000]}"  # Limit each doc to 2000 chars
            for doc in documents
        ])
        
        # Generate report using OpenAI
        report = self.openai_manager.generate_report(combined_content)
        
        # Add metadata
        report['metadata'] = {
            'generated_at': datetime.datetime.now().isoformat(),
            'source_documents': [doc['filename'] for doc in documents],
            'total_documents': len(documents)
        }
        
        return report
    
    def generate_from_topic(self, topic: str) -> Dict:
        report = self.openai_manager.generate_report(None, topic)
        
        # Add metadata
        report['metadata'] = {
            'generated_at': datetime.datetime.now().isoformat(),
            'source_type': 'topic_query',
            'topic': topic
        }
        
        return report
    
    def _create_empty_report(self) -> Dict:
        """Create an empty report structure"""
        return {
            "Executive Summary": "No documents provided for analysis.",
            "metadata": {
                'generated_at': datetime.datetime.now().isoformat(),
                'source_type': 'empty',
                'total_documents': 0
            }
        }


class PDFGenerator:
    """Generates formatted PDF reports"""
    
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self.create_custom_styles()
    
    def create_custom_styles(self):
        self.styles.add(ParagraphStyle(
            name='CustomTitle',
            parent=self.styles['Heading1'],
            fontSize=24,
            spaceAfter=30,
            alignment=TA_CENTER
        ))
        
        self.styles.add(ParagraphStyle(
            name='SectionHeading',
            parent=self.styles['Heading2'],
            fontSize=16,
            spaceBefore=20,
            spaceAfter=10,
            textColor=colors.HexColor('#2E3440')
        ))
        
        self.styles.add(ParagraphStyle(
            name='CustomBody',
            parent=self.styles['Normal'],
            fontSize=11,
            alignment=TA_JUSTIFY,
            spaceAfter=12
        ))
    
    def create_pdf(self, report: Dict, output_path: str):
        doc = SimpleDocTemplate(
            output_path,
            pagesize=letter,
            rightMargin=0.75*inch,
            leftMargin=0.75*inch,
            topMargin=1*inch,
            bottomMargin=0.75*inch
        )
        
        story = []
        
        # Title page
        story.append(Paragraph("Engineering Meeting Brief", self.styles['CustomTitle']))
        story.append(Spacer(1, 0.5*inch))
        story.append(Paragraph(
            f"Generated on: {datetime.datetime.now().strftime('%B %d, %Y %H:%M')}",
            self.styles['Normal']
        ))
        
        if 'metadata' in report:
            if 'source_documents' in report['metadata']:
                story.append(Paragraph(
                    f"Based on {report['metadata']['total_documents']} documents",
                    self.styles['Normal']
                ))
            elif 'topic' in report['metadata']:
                story.append(Paragraph(
                    f"Topic: {report['metadata']['topic']}",
                    self.styles['Normal']
                ))
        
        story.append(PageBreak())
        
        # Add each section from the report
        section_order = [
            'Executive Summary',
            'Topic Overview',
            'Technology Stack Analysis',
            'Architecture Overview',
            'Advantages and Benefits',
            'Limitations and Challenges',
            'Alternative Solutions',
            'Competitive Analysis',
            'Key Recommendations',
            'Action Items and Next Steps'
        ]
        
        for section in section_order:
            if section in report and report[section]:
                story.append(Paragraph(section, self.styles['SectionHeading']))
                content = report[section]
                
                if isinstance(content, list):
                    for item in content:
                        story.append(Paragraph(f"• {str(item)}", self.styles['CustomBody']))
                else:
                    paragraphs = str(content).split('\n\n')
                    for paragraph in paragraphs:
                        if paragraph.strip():
                            story.append(Paragraph(paragraph, self.styles['CustomBody']))
                
                story.append(Spacer(1, 0.25*inch))
        
        # Build PDF
        try:
            doc.build(story)
        except Exception as e:
            console.print(f"[red]Error creating PDF: {e}[/red]")
            raise


def main():
    load_dotenv()
    
    parser = argparse.ArgumentParser(
        description="Meeting Preparation Tool - Generate technical briefings from documents or topics",
        formatter_class=argparse.RawDescriptionHelpFormatter,
        epilog="""
Examples:
  meeting-prep --path ./docs --output tech_brief.pdf
  meeting-prep --topic "kubernetes vs docker swarm" --output comparison.pdf
  meeting-prep  # Interactive mode
        """
    )
    parser.add_argument("--path", "-p", type=str, help="Path to documents folder")
    parser.add_argument("--topic", "-t", type=str, help="Topic for direct query (no documents)")
    parser.add_argument("--output", "-o", type=str, default="meeting_brief.pdf", help="Output PDF filename")
    parser.add_argument("--verbose", "-v", action="store_true", help="Verbose output")
    
    args = parser.parse_args()
    
    if args.path and args.topic:
        console.print("[red]Error:[/red] Please provide either --path or --topic, not both.")
        sys.exit(1)
    
    if not args.path and not args.topic:
        # Interactive mode
        console.print("[bold blue]Meeting Prep Tool[/bold blue]")
        console.print("Generate technical briefings for your meetings\n")
        
        choice = Prompt.ask("Choose input method", choices=["documents", "topic"])
        if choice == "documents":
            args.path = Prompt.ask("Enter path to documents folder")
        else:
            args.topic = Prompt.ask("Enter the topic for briefing")
    
    try:
        with Progress(
            SpinnerColumn(),
            TextColumn("[progress.description]{task.description}"),
            console=console
        ) as progress:
            if args.path:
                # Document-based processing
                task = progress.add_task("Processing documents...", total=None)
                processor = DocumentProcessor(Path(args.path), verbose=args.verbose)
                documents = processor.process_all()
                
                if not documents:
                    console.print("[yellow]No documents were processed. Please check your input path.[/yellow]")
                    sys.exit(1)
                
                progress.update(task, description="Generating report...")
                generator = ReportGenerator()
                report = generator.generate_from_content(documents)
            else:
                # Topic-based processing
                task = progress.add_task("Generating report from topic...", total=None)
                generator = ReportGenerator()
                report = generator.generate_from_topic(args.topic)
            
            # Generate PDF
            progress.update(task, description="Creating PDF...")
            pdf_gen = PDFGenerator()
            pdf_gen.create_pdf(report, args.output)
        
        console.print(f"[green]✓[/green] Report generated: {args.output}")
        
    except FileNotFoundError as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)
    except ValueError as e:
        console.print(f"[red]Error:[/red] {e}")
        sys.exit(1)
    except Exception as e:
        console.print(f"[red]Unexpected error:[/red] {e}")
        if args.verbose:
            import traceback
            console.print(traceback.format_exc())
        sys.exit(1)


if __name__ == "__main__":
    load_dotenv()
    
    # Check if API key is set
    if not os.getenv('OPENAI_API_KEY'):
        console.print("[red]Error:[/red] OPENAI_API_KEY environment variable not set")
        console.print("Please set it using: export OPENAI_API_KEY='your-api-key'")
        sys.exit(1)
    main()
