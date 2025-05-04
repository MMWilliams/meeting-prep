from pathlib import Path
import PyPDF2
from PIL import Image
import pytesseract
from pptx import Presentation
import docx
import chardet

class ContentExtractor:
    def extract(self, file_path: Path) -> str:
        """Extract text content from various file formats"""
        suffix = file_path.suffix.lower()
        
        if suffix == '.pdf':
            return self.extract_pdf(file_path)
        elif suffix == '.txt':
            return self.extract_text(file_path)
        elif suffix in ['.png', '.jpg', '.jpeg']:
            return self.extract_image(file_path)
        elif suffix == '.pptx':
            return self.extract_pptx(file_path)
        elif suffix == '.docx':
            return self.extract_docx(file_path)
        else:
            raise ValueError(f"Unsupported file type: {suffix}")
    
    def extract_pdf(self, file_path: Path) -> str:
        """Extract text from PDF files"""
        text = ""
        with open(file_path, 'rb') as file:
            pdf_reader = PyPDF2.PdfReader(file)
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
        return text
    
    def extract_text(self, file_path: Path) -> str:
        """Extract text from text files with encoding detection"""
        with open(file_path, 'rb') as file:
            raw_data = file.read()
            result = chardet.detect(raw_data)
            encoding = result['encoding'] or 'utf-8'
        
        with open(file_path, 'r', encoding=encoding) as file:
            return file.read()
    
    def extract_image(self, file_path: Path) -> str:
        """Extract text from images using OCR"""
        image = Image.open(file_path)
        return pytesseract.image_to_string(image)
    
    def extract_pptx(self, file_path: Path) -> str:
        """Extract text from PowerPoint presentations"""
        prs = Presentation(file_path)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    
    def extract_docx(self, file_path: Path) -> str:
        """Extract text from Word documents"""
        doc = docx.Document(file_path)
        return "\n".join([paragraph.text for paragraph in doc.paragraphs])
